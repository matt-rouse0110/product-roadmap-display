import random
import requests
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN

def disp_message():
   print("Hello!!")
def executeRoadmapCreation(username, APIKey, project, fields):
   
   print(project)
   print(fields)
   url = 'https://gentrack.atlassian.net/rest/api/3/search?'
   jql='jql=project%3D'+project+'+and+issuetype%3DEpic&fields='+fields
   print(url+jql)
   myResponse = requests.get(url+jql,auth=(username,APIKey))
   # f = open(r'C:\Users\mattr\OneDrive\Documents\pythonTesting\dummyData.json', 'r')
   # jsonData = json.load(f)
   jsonData = myResponse.json()
   print(json.dumps(jsonData,indent=3))
   
   labels = ['Now','Next','Later']
   # print(jsonData["issues"])

   prs = Presentation()
   title_only_slide_layout = prs.slide_layouts[5]
   slide_width = prs.slide_width
   slide_height = prs.slide_height
   slide = prs.slides.add_slide(title_only_slide_layout)
   shapes = slide.shapes

   shapes.title.text = 'Product Roadmap'

   left = slide_width/3-(slide_width/3)/2
   top = Inches(3.0)
   width = Inches(1.75)
   height = Inches(1.0)

   shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
   shape.text = 'Now'

   Nextleft = left + (slide_width)/3
   shape = shapes.add_shape(MSO_SHAPE.PENTAGON, Nextleft, top, width, height)
   shape.text = 'Next'

   Laterleft = Nextleft + (slide_width)/3
   print(left)
   print(Nextleft)
   print(Laterleft)
   shape = shapes.add_shape(MSO_SHAPE.PENTAGON, Laterleft, top, width, height)
   shape.text = 'Later'
   
   top = top + height + Inches(0.1)
   Nexttop = top + height + Inches(0.1)
   Latertop = top + height + Inches(0.1)
   width = Inches(1.75)  # chevrons need more width for visual balance
   for x in jsonData['issues']:
      label = random.choice(labels)
      # print(x['key']+','+label)
      if label == 'Now':
         shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
         top = top + height  + Inches(0.1)
      elif label == 'Next':
         shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, Nextleft, Nexttop, width, height)
         Nexttop = Nexttop + height  + Inches(0.1)
      elif label == 'Later':
         shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, Laterleft, Latertop, width, height)
         Latertop = Latertop + height  + Inches(0.1)
      shape.text = x['key']
      shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
      shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
      


   prs.save('test.pptx')